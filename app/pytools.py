'''
Copyright 2023, Thomas Wright

LICENSE NOTICE
This file is part of The Pep Genie.
The Pep Genie is free software: 
you can redistribute it and/or modify it under the terms of the 
GNU General Public License as published by the Free Software Foundation, 
either version 3 of the License, or (at your option) any later version.
The Pep Genie is distributed in the hope that it will be useful, 
but WITHOUT ANY WARRANTY; without even the implied warranty of 
MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE. 
See the GNU General Public License for more details.
You should have received a copy of the GNU General Public License along with The Pep Genie.
If not, see <https://www.gnu.org/licenses/>. 


'''



from asyncore import file_dispatcher
import cv2
import numpy as np
import ast
import os
import zipfile
from PIL import Image, ImageStat, ImageDraw, ImageOps
import matplotlib.pyplot as plt

from django.conf import settings
media_root = settings.MEDIA_ROOT

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm


def strip_number(input_string, first_col):
    
    input_string = input_string.replace(' ','')
    input_list = input_string.split(',')

    letter_ref = ['7','A', 'B', 'C', 'D', 'E', 'F', 'G', 'H', 'I', 'J', 'K', 'L', 'M', 'N', 'O', 'P', 'Q', 'R', 'S', 'T', 'U', 'V', 'W', 'X', 'Y', 'Z']

    new_strip = []
    strip_list=[]
    for strip in input_list:
 
        #convert letter to row number
        if letter_ref.index(strip[0]) < 10:
            new_strip.append('0')
            new_strip.append(letter_ref.index(strip[0]))
        else:
            new_strip.append(letter_ref.index(strip[0]))
        #add column number
        new_strip.append(strip[1:3])
        new_strip.append('-')

        #convert letter 2 to row number
        if letter_ref.index(strip[4]) < 10:
            new_strip.append('0')
            new_strip.append(letter_ref.index(strip[4]))
        else:
            new_strip.append(letter_ref.index(strip[4]))
        #add column number
        new_strip.append(strip[5:7])

        #add output to list
        strip_list.append(''.join(map(str, new_strip)))
        new_strip = []


    new_item = 0
    new_row = []
    result=[]
    for i in strip_list:
        #add row number
        new_item += int(i[0:2])
        #add col number
        new_item += (int(i[2:4])-first_col)*20
        new_row.append(new_item)
        new_item = 0

        new_item += int(i[5:7])
        new_item += (int(i[7:9])-first_col)*20
        new_row.append(new_item)
        new_item = 0
        result.append(new_row)
        new_row = []
    #for i in result:
        #print(i[0], i[1])

    return result




def generate_grid(rows,cols,path):
    grid_unit = cv2.imread(path, -1)
    my_list = []
    for i in range(int(rows)):
        my_list.append(grid_unit)
    grid_column = cv2.vconcat(my_list)
    my_list=[]
    for i in range(int(cols)):
        my_list.append(grid_column)
    full_grid = cv2.hconcat(my_list)
    #cv2.imwrite(out_path, full_grid)
    return full_grid







def slice_to_list(my_image, rows, cols):
    #variables
    #my_image = Image.open(image_path)
    old_width, old_height = my_image.size #was issue opening image from model, but 
                                        # now I am supplying the PIL file

    #resize cropped image
    if old_height <= 400:
        new_width = int( (400 / old_height) * old_width )
        my_image = my_image.resize((new_width, 400))

    #resized image variables
    width, height = my_image.size
    slice_height = height / rows
    slice_width = width / cols
    upper = 0
    left = 0

    #cut cropped image into squares/slices
    count = 1
    slice_list = []
    for col in range(cols):
        right = (col+1)* slice_width
        for row in range(rows):
            if count == rows:
                lower = height
            else:
                lower = int(count * slice_height)
            bbox = (left, upper, right, lower)
            working_slice = my_image.crop(bbox)
            upper += int(slice_height)
            slice_list.append(working_slice)
            my_width, my_height = working_slice.size
            #print(count, my_height, lower, upper)
            count += 1
        left += slice_width
        upper = 0
        count = 1
    return slice_list




#Crop image to circle and measure

#mask variables
#width = 200
def make_mask(width, spot_diameter):
    frac = width * ( (1- (spot_diameter / 100)) / 2)
    print('spot diameter and width frac')
    print(spot_diameter)
    print(frac)


    #create mask image
    bbox = [frac, frac, width-frac, width-frac]
    mask = Image.new('L', (width, width),0)
    draw = ImageDraw.Draw(mask)
    draw.pieslice(bbox,0,360,fill=255)

    
    return mask





def circle_crop(im, mask):
    #crop image to mask
    im = ImageOps.fit(im, mask.size, centering=(0.5, 0.5))
    cropped = im.convert('L')
    #add alpha layer to image
    cropped.putalpha(mask)
    #output is RGBA

    #new method
    #image_array = np.array(im)
    #mask_array = np.array(mask)
    #cropped = Image.fromarray(np.dstack((image_array,mask_array)))

    #test
    #cropped = Image.fromarray(image_array)


    return cropped



def measure_darkness(im):
    stat = ImageStat.Stat(im)
    return stat.mean[0]





def stripper2(image_path, graph_type_list, rows, cols, strips):
    #variables
    my_image = Image.open(image_path)
    old_width, old_height = my_image.size

    #resize cropped image
    new_width = int( (1000 / old_height) * old_width )
    my_image = my_image.resize((new_width, 1000))

    #resized image variables
    width, height = my_image.size
    slice_height = height / rows
    slice_width = width / cols
    upper = 0
    left = 0

    #cut cropped image into squares/slices
    count = 1
    slice_list = []
    for col in range(cols):
        right = (col+1)* slice_width
        for row in range(rows):
            if count == rows:
                lower = height
            else:
                lower = int(count * slice_height)
            bbox = (left, upper, right, lower)
            working_slice = my_image.crop(bbox)
            upper += int(slice_height)
            slice_list.append(working_slice)
            my_width, my_height = working_slice.size
            #print(count, my_height, lower, upper)
            count += 1
        left += slice_width
        upper = 0
        count = 1

    # convert input string to number tuple
    output_strips = []
    for number_pair in strips:
        new_strip = []
        new_strip.append(number_pair[0]-1)
        new_strip.append(number_pair[1])
        image_number = new_strip[1]-new_strip[0]

        #Make strip canvas and join squares/slices
        width, height = slice_list[0].size

        new_image = Image.new('RGB', (width, height * image_number))

        counter = 0
        index_counter = new_strip[0]+1
        for image in slice_list[new_strip[0]:new_strip[1]]:
            new_image.paste(image,(0,counter*height))
            if ((index_counter-1)/20).is_integer() and counter != 0:
                xy = (0,counter*height,width,counter*height)
                draw = ImageDraw.Draw(new_image)
                draw.line(xy, fill='black', width=3)
            index_counter += 1
            counter +=1
        if graph_type_list[strips.index(number_pair)] == 'v':
            new_image = new_image.transpose(Image.ROTATE_90)
            new_image = new_image.resize((height*image_number,width)) 
        output_strips.append(new_image) 
        width, height = output_strips[0].size 


    return output_strips







### MAKE GRAPHS ###---------------------------------------


def h_bar(my_x,my_heights,my_dir,my_filename,ylim=False):
    plt.switch_backend('Agg') 
    plt.clf()
    plt.barh(my_x,my_heights)

    #invert y
    ax = plt.gca() #get current axis
    ax.invert_yaxis()

    #hide y ticks
    ax.yaxis.set_ticklabels([])
    plt.yticks([], [])
    

    #add threshold
    if my_x[0] != 0:
        plt.axvline(x=1,color='blue',lw=1, ls='--')

    #hide frame/border
    ax.spines['bottom'].set_visible(False)
    ax.spines['right'].set_visible(False)

    #ax.set_xlabel('Relative Luminescence',fontsize=14)

    #x axis on top
    ax.xaxis.tick_top()


    plt.margins(y=0)

    #set axis if disc array
    if ylim != False:
        plt.xlim(0,ylim)

    #plt.show()
    fig = plt.gcf()
    spot_no = len(my_x)
    width = (7.22/20)*spot_no
    fig.set_size_inches(5, width) #8.37/20 now 7.2/20 per seq, width 6

    my_path = os.path.join(my_dir,my_filename)
    fig.savefig(my_path, bbox_inches='tight',dpi=100)






def v_bar(my_x,my_heights,my_dir,my_filename,ylim=False):
    plt.switch_backend('Agg') 
    plt.clf()
    plt.bar(my_x,my_heights)

    ax = plt.gca() #get current axis    

    #add threshold
    if my_x[0] != 0:
        plt.axhline(y=1,color='blue',lw=1, ls='--')

    #hide frame/border
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    #hide y ticks
    ax.xaxis.set_ticklabels([])
    plt.xticks([], [])

    #ax.set_xlabel('Relative Luminescence',fontsize=14)

    plt.margins(x=0)
    #plt.show()

    #set axis if disc array
    if ylim != False:
        plt.ylim(0,ylim)

    fig = plt.gcf()
    spot_no = len(my_x)
    width = (7.22/20)*spot_no
    fig.set_size_inches(width,4) #8.37/20 now 7.2/20 per seq, width 6
    my_path = os.path.join(my_dir,my_filename)
    fig.savefig(my_path, bbox_inches='tight',dpi=100)



def file_compress(inp_file_names, filenames, out_zip_file):
    compression = zipfile.ZIP_DEFLATED
    zf = zipfile.ZipFile(out_zip_file, mode="w")
    try:
        for i in range(len(inp_file_names)):
    
            zf.write(inp_file_names[i], filenames[i], compress_type=compression)

    except FileNotFoundError as e:
        print('compression error')
    finally:
        zf.close()









def make_pptx(test_list, graph_type_list,strip_request,user_dir,user_id,vis_bool):
    prs = Presentation()
    
    prs.slide_width = Cm(33.9)
    prs.slide_height = Cm(19.1)
    
    for i in range(len(test_list)):
        #get number of slices for height calculation
        new_strip = []
        new_strip.append(strip_request[i][0]-1)
        new_strip.append(strip_request[i][1])
        image_number = new_strip[1]-new_strip[0]
        long = Cm(((20.3*0.7)/20)*image_number)
        short = Cm(0.7)
        #account for vertical vs horizontal
        if graph_type_list[i] == 'h': #functional
            width = short
            height = long
            left = Cm(8) #one cm before
            left_mock = Cm(7)
            left_vis = Cm(6)
            top = Cm(3)
            top_mock = Cm(3)
            top_vis = Cm(3)
            
            text_left = Cm(6.75)
            text_top = Cm(2)
            text_width = Cm(1)
            text_height = Cm(1)
            text_align = PP_ALIGN.LEFT
            text_rotation = -90.0
            text_left_factor = Cm(1.1)
            text_top_factor = Cm(0)

        elif graph_type_list[i] == 'v': #vertical bars
            width = long
            height = short
            left = Cm(9.5) #one cm before
            left_mock = Cm(9.5)
            left_vis = Cm(9.5)
            top = Cm(15) #
            top_mock = Cm(16) #
            top_vis = Cm(17) #
 
            text_left = Cm(2.5)
            text_top = Cm(14.9)
            text_width = Cm(7)
            text_height = Cm(1)
            text_align = PP_ALIGN.RIGHT
            text_rotation = 0
            text_left_factor = Cm(0)
            text_top_factor = Cm(0.9)

        slide = prs.slides.add_slide(prs.slide_layouts[5])

        my_file_name = 'strip_' + str(i) + '.png'
        my_path = os.path.join(user_dir, my_file_name)
        print(width,height)
        img=slide.shapes.add_picture(my_path,left,top,width,height)

        my_file_name_mock = 'mock_strip_' + str(i) + '.png'
        my_path_mock = os.path.join(user_dir, my_file_name_mock)
        img = slide.shapes.add_picture(my_path_mock,left_mock,top_mock,width,height)

        if vis_bool == True:
            my_file_name_vis = 'vis_strip_' + str(i) + '.png'
            my_path_vis = os.path.join(user_dir, my_file_name_vis)
            img = slide.shapes.add_picture(my_path_vis,left_vis,top_vis,width,height)


        #add labels for strips
        txBox = slide.shapes.add_textbox(text_left, text_top+text_top_factor, text_width, text_height)
        tf = txBox.text_frame
        tf.text = "Control"
        txBox.rotation = text_rotation
        txBox.text_frame.paragraphs[0].alignment = text_align
        
        txBox2 = slide.shapes.add_textbox(text_left+text_left_factor, text_top, text_width, text_height)
        tf2 = txBox2.text_frame
        tf2.text = "Overlay"
        txBox2.rotation = text_rotation
        txBox2.text_frame.paragraphs[0].alignment = text_align

        if vis_bool == True:
            txBox3 = slide.shapes.add_textbox(text_left+text_left_factor*(-1), text_top+text_top_factor*(2), text_width, text_height)
            tf3 = txBox3.text_frame
            tf3.text = "UV/Vis"
            txBox3.rotation = text_rotation
            txBox3.text_frame.paragraphs[0].alignment = text_align
    

    my_filename = 'presentation_'+user_id+'.pptx'
    prs.save(os.path.join(user_dir,my_filename))

#second function customised for sd



