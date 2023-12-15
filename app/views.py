'''
Copyright 2023, Thomas Wright

LICENSE NOTICE
This file is part of The Pep Genie.
The Pep Genie is free software: 
you can redistribute it and/or modify it under the terms of the 
GNU General Public License as published by the Free Software Foundation, 
either version 3 of the License, or (at your option) any later version.
The Pep Genie is distributed in the hope that it will be useful, 
but WITHOUT ANY WARRANTY; without even the implied warranty of 
MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE. 
See the GNU General Public License for more details.
You should have received a copy of the GNU General Public License along with The Pep Genie.
If not, see <https://www.gnu.org/licenses/>. 


'''



from django.shortcuts import render, redirect
from datetime import datetime
from .models import gridModel
from .forms import gridModelForm
from .pytools import *
from .seq_studies import *
from django.conf import settings
import os
import cv2
import csv
from pptx import Presentation
from pptx.util import Cm
max_ = max;
min_ = min;
from django.core.files.base import ContentFile
from PIL import Image
import numpy as np
import ast
from django.http import HttpResponse, Http404, HttpResponseServerError

MEDIA_ROOT = settings.MEDIA_ROOT




def index(request):
    print(np.__version__)
    print(cv2.__version__)
    return render(request, 'app/index.html',{})

# Create your views here.
def request(request):
    #make new instance and set user_id as four digit number
    #generate id
    this_now = datetime.now()
    hash = this_now.strftime("%Y-%m-%d-%H-%M-%S")
    my_model = gridModel(user_id=hash)

    #save default image to model for vis so optional
    my_image = Image.open(os.path.join(settings.BASE_DIR, 'app','static','app','img','dummy.png'))
    my_image = np.array(my_image)
    ret, buf = cv2.imencode('.png', my_image) 
    content = ContentFile(buf.tobytes())
    my_model.img_vis.save('EMPTY.jpg', content)

    #also save for test and vis
    my_image = Image.open(os.path.join(settings.BASE_DIR, 'app','static','app','img','demo_array.jpeg'))
    my_image = np.array(my_image)
    ret, buf = cv2.imencode('.png', my_image) 
    content = ContentFile(buf.tobytes())
    my_model.img_test.save('DEMO.jpg', content)
    my_model.img_con.save('DEMO.jpg', content)



    my_model.save() 

    #load form of that existing instance
    grid_form = gridModelForm(instance = my_model)
    #return render(request, 'lab/sd_upload.html',{'form':sd_form})
    return render(request, 'app/request.html', {'form':grid_form})

'''
def request(request):
    return render(request, 'app/request.html', {})

def grid(request):
    return render(request, 'app/grid.html', {})

def result(request):
    return render(request, 'app/result.html', {})

'''

def validate(request):
    my_model = gridModel.objects.all().order_by('-id')[0]

    #Validate form and save model
    if request.method == 'POST':
        form = gridModelForm(request.POST, request.FILES,  instance=my_model)
        if form.is_valid():
            form.save()
    else:
        args={}
        form = gridModelForm()
        args['form'] = form
        return render(request, 'app/request.html',args)

    #convert strip_request to usable list
    my_model.strip_request_string = my_model.strip_request
    my_model.strip_request = strip_number(my_model.strip_request,my_model.first_col)
    my_model.save()

    #format graph_types
    my_string = my_model.graph_types
    my_model.graph_types_string = my_string
    my_string.replace(' ','')
    my_list = my_string.split(',')

    my_model.graph_types = my_list
    #format graph types if discovery array - y axes need to be kept the same if discovery array
    if 'd' in my_list:
        new_list = []
        for i in range(len(my_model.strip_request)):
            new_list.append('d')
        my_model.graph_types = new_list
    my_model.save()
    print(my_model.graph_types)
    print('^GRAPH TYPES')

    #make context for graph types template
    my_list = my_model.strip_request_string.split(',')
    print(my_model.strip_request)
    print('^STRIP REQUEST')
    print(my_list)
    print('^MY LIST')


    #generate grid file
    rows = 20
    cols = my_model.cols
    grid_unit_path = os.path.join(settings.BASE_DIR, 'app','static','app','img','grid-unit.png')
    full_grid = generate_grid(rows,cols,grid_unit_path)
    #save full grid to model
    ret, buf = cv2.imencode('.png', full_grid) 
    content = ContentFile(buf.tobytes())
    my_model.full_grid.save('full_grid.png', content)   

    return redirect('app:graph_types') 


def graph_types(request):
    my_model = gridModel.objects.all().order_by('-id')[0]
    my_list = my_model.strip_request_string
    my_list = my_list.split(',')

    context = {'my_list':my_list}
    return render(request, 'app/graph_types.html',context)






def grid_test(request):

    my_model = gridModel.objects.all().order_by('-id')[0]

    graph_types = request.POST.get('hiddenField')
    graph_types.replace(' ','')
    graph_types = graph_types.split(',')

    if 'd' in graph_types:
        new_list = []
        for i in range(len(my_model.strip_request)):
            new_list.append('d')
        graph_types = new_list
    my_model.graph_types = graph_types
    my_model.save()
    header = 'Test Array'
    context = {'my_model':my_model,'img':my_model.img_test,'header':header}
    return render(request, 'app/grid_test.html',context) 


def grid_con(request):    
    my_model = gridModel.objects.all().order_by('-id')[0]

    x1 = float(request.POST.get('coordX1'))
    x2 = float(request.POST.get('coordX2'))
    y1 = float(request.POST.get('coordY1'))
    y2 = float(request.POST.get('coordY2'))

    #check if test grid was not drawn. C set to 1 before template. 
    if x1 == x2:
        print('grid not drawn after test array. Sending back')
        header = 'Test Array'
        context = {'my_model':my_model,'img':my_model.img_test,'header':header}
        return render(request, 'app/grid_test.html',context) 
    else:
        print('grid drawn')
        print('x',x1)
        print('x',x2)
    canvas_width = float(request.POST.get('canvas_width'))
    canvas_height = float(request.POST.get('canvas_height'))
    my_coords = [x1,y1,x2,y2]
    my_model.img_test_coords = my_coords

    my_model.save()

    header = 'Control Array'
    context = {'my_model':my_model,'img':my_model.img_con,'header':header,
                'x1':''.join([str(round(x1*canvas_width)),'px']),
                'y1':''.join([str(round(y1*canvas_height)),'px']),
                'my_width':''.join([str(round((x2-x1)*canvas_width)),'px']),
                'my_height':''.join([str(round((y2-y1)*canvas_height)),'px'])
                }

    return render(request, 'app/grid_con.html',context) 


def grid_vis(request):
    my_model = gridModel.objects.all().order_by('-id')[0]
    #save coords
    x1 = float(request.POST.get('coordX1'))
    y1 = float(request.POST.get('coordY1'))
    x2 = float(request.POST.get('coordX2'))
    y2 = float(request.POST.get('coordY2'))
    my_coords = [x1,y1,x2,y2]
    my_model.img_con_coords = my_coords
    my_model.save()

    #If vis is EMPTY then skip to download
    my_string = my_model.img_vis.name
    last_slash_index = my_string.rindex('/')
    if my_string[last_slash_index + 1:] =='EMPTY.jpg':
        #skip vis
        return redirect('app:grid_result')
    else:
        #return to grid for uv
        header = 'UV/Vis Control Array'
        context = {'my_model':my_model,'img':my_model.img_vis,'header':header}
        return render(request, 'app/grid_vis.html',context) 


def grid_vis_after(request):
    my_model = gridModel.objects.all().order_by('-id')[0]
    x1 = float(request.POST.get('coordX1'))
    y1 = float(request.POST.get('coordY1'))
    x2 = float(request.POST.get('coordX2'))
    y2 = float(request.POST.get('coordY2'))
    #test if grid was drawn
    if x1 == x2:
        header = 'UV/Vis Control Array'
        context = {'my_model':my_model,'img':my_model.img_vis,'header':header}
        return render(request, 'app/grid_vis.html',context) 
    my_coords = [x1,y1,x2,y2]
    my_model.img_vis_coords = my_coords
    my_model.save()
    
    return redirect('app:grid_result')










def download(request):
    my_model = gridModel.objects.all().order_by('-id')[0]
    return render(request, 'app/download.html', {})












def grid_result(request):
    my_model = gridModel.objects.all().order_by('-id')[0]

    #crop img1
    test_coords = ast.literal_eval(my_model.img_test_coords)
    my_x1, my_y1, my_x2, my_y2 = test_coords[0],test_coords[1],test_coords[2],test_coords[3]
    try:
        my_url = os.path.join(MEDIA_ROOT,my_model.img_test.name)
        image = Image.open(my_url)
    except:
        my_url = os.path.normpath(MEDIA_ROOT+my_model.img_test.name)
        image = Image.open(my_url)
    width, height = image.size
    left = float(my_x1)*float(width)
    upper = float(my_y1)*float(height)
    right = float(my_x2)*float(width)
    lower = float(my_y2)*float(height)
    img_test_cropped = image.crop((left, upper, right, lower))
    #save cropped image
    #convert PIL to openCV
    open_cv_image = np.array(img_test_cropped) 
    # Convert RGB to BGR 
    open_cv_image = cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2BGR)
    #save cropped_image to model
    ret, buf = cv2.imencode('.png', open_cv_image) 
    content = ContentFile(buf.tobytes())
    my_model.img_test_crop.save('img_test_cropped.png', content)
    #my_model.save()  
     
    
    #crop img con
    test_coords = ast.literal_eval(my_model.img_con_coords)
    my_x1, my_y1, my_x2, my_y2 = test_coords[0],test_coords[1],test_coords[2],test_coords[3]
    try:
        my_url = os.path.join(MEDIA_ROOT,my_model.img_con.name)
        image = Image.open(my_url)
    except:
        my_url = os.path.normpath(MEDIA_ROOT+my_model.img_con.name)
        image = Image.open(my_url)
    width, height = image.size
    left = float(my_x1)*float(width)
    upper = float(my_y1)*float(height)
    right = float(my_x2)*float(width)
    lower = float(my_y2)*float(height)
    img_con_cropped = image.crop((left, upper, right, lower))
    #save cropped image
    #convert PIL to openCV
    open_cv_image = np.array(img_con_cropped) 
    # Convert RGB to BGR 
    open_cv_image = cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2BGR)
    #save cropped_image to model
    ret, buf = cv2.imencode('.png', open_cv_image) 
    content = ContentFile(buf.tobytes())
    my_model.img_con_crop.save('img_con_cropped.png', content)
    #my_model.save() 



    #if vis is not empty then crop vis too
    my_string = my_model.img_vis.name
    last_slash_index = my_string.rindex('/')
    if my_string[last_slash_index + 1:] !='EMPTY.jpg':
        
        #crop img vis
        test_coords = ast.literal_eval(my_model.img_vis_coords)
        my_x1, my_y1, my_x2, my_y2 = test_coords[0],test_coords[1],test_coords[2],test_coords[3]
        try:
            my_url = os.path.join(MEDIA_ROOT,my_model.img_vis.name)
            image = Image.open(my_url)
        except:
            my_url = os.path.normpath(MEDIA_ROOT+my_model.img_vis.name)
            image = Image.open(my_url)
        width, height = image.size
        left = float(my_x1)*float(width)
        upper = float(my_y1)*float(height)
        right = float(my_x2)*float(width)
        lower = float(my_y2)*float(height)
        img_vis_cropped = image.crop((left, upper, right, lower))
        #save cropped image
        #convert PIL to openCV
        open_cv_image = np.array(img_vis_cropped) 
        # Convert RGB to BGR 
        open_cv_image = cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2BGR)
        #save cropped_image to model
        ret, buf = cv2.imencode('.png', open_cv_image) 
        content = ContentFile(buf.tobytes())
        my_model.img_vis_crop.save('img_vis_cropped.png', content)
        #my_model.save() 


    ### Measure Darkness and save to csv files ###

    #get darkness list of img1
    rows = 20
    cols = my_model.cols
    slice_list = slice_to_list(img_test_cropped, rows, cols)

    #NEED TO MELT THE SLICE LIST, AS CSV WRITES IN ROWS, NOT COLS
    format_list = []
    c=0
    r=0
    for i in range(rows):
        for i in range(cols):
            format_list.append(slice_list[c+r])
            c+=rows
        r+=1
        c=0

    #crop each image in list and measure darkness
    mask = make_mask(200, my_model.measure_diameter)


    
    #measure darknesses
    darkness_list_img1 = []
    my_row = []

    c = 0
    for row in range(rows):
        for col in range(cols):
            cropped = circle_crop(format_list[c],mask)
            my_row.append(measure_darkness(cropped))
            c += 1
        darkness_list_img1.append(my_row)
        my_row = []


    #get darkness list of img2
    rows = 20
    cols = my_model.cols
    slice_list = slice_to_list(img_con_cropped, rows, cols)

    #NEED TO MELT THE SLICE LIST, AS CSV WRITES IN ROWS, NOT COLS
    format_list = []
    c=0
    r=0
    for i in range(rows):
        for i in range(cols):
            format_list.append(slice_list[c+r])
            c+=rows
        r+=1
        c=0

    #crop each image in list and measure darkness
    mask = make_mask(200, my_model.measure_diameter)
    
    #measure darknesses
    darkness_list_img2 = []
    my_row = []

    c = 0
    for row in range(rows):
        for col in range(cols):
            cropped = circle_crop(format_list[c],mask)
            my_row.append(measure_darkness(cropped))
            c += 1
        darkness_list_img2.append(my_row)
        my_row = []


    #calculate min, max and delta for darkness transformation
    darkness_list_values = [] 

    for row in range(rows):
        for col in range(cols):
            darkness_list_values.append(darkness_list_img1[row][col])
            darkness_list_values.append(darkness_list_img2[row][col])

    my_min = min_(darkness_list_values)
    my_max = max_(darkness_list_values)
    my_delta =  my_max-my_min


    #transform data so that values increase from light to dark
    darkness_list_img_1_trn = []
    my_row = []
    for row in range(rows):
        for col in range(cols):
            #invert value while retaining distance from 0
            my_row.append(
                ((1-((darkness_list_img1[row][col]-my_min)/my_delta))*my_delta)+my_min
            )
        darkness_list_img_1_trn.append(my_row)
        my_row = []

    #export for making graphs
    #if norm = false, then this will be fed to graph
    my_model.dens_list = darkness_list_img_1_trn
    my_model.save()

    darkness_list_img_2_trn = []
    my_row = []
    for row in range(rows):
        for col in range(cols):
            #invert value while retaining distance from 0
            my_row.append(
                ((1-((darkness_list_img2[row][col]-my_min)/my_delta))*my_delta)+my_min
            )
        darkness_list_img_2_trn.append(my_row)
        my_row = []


    #save test image darknesses
    path_name = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
    path_name = os.path.join(path_name,my_model.user_id) 
    path_name = os.path.join(path_name,'img_test.csv')

    with open(path_name, 'w', newline='') as csvfile:
        writer = csv.writer(csvfile)
        writer.writerows(darkness_list_img_1_trn)
    #save control image darknesses

    path_name = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
    path_name = os.path.join(path_name,my_model.user_id) 
    path_name = os.path.join(path_name,'img_con.csv')

    with open(path_name, 'w', newline='') as csvfile:
        writer = csv.writer(csvfile)
        writer.writerows(darkness_list_img_2_trn)

    #background subtract 

    my_row = []
    darkness_list_background_subtracted = []
    for row in range(rows):
        for col in range(cols):
            my_row.append(darkness_list_img_1_trn[row][col] - darkness_list_img_2_trn[row][col])
        darkness_list_background_subtracted.append(my_row)
        my_row = []

    #save normalised data
    path_name = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
    path_name = os.path.join(path_name,my_model.user_id) 
    path_name = os.path.join(path_name,'background_subtracted_data.csv')

    with open(path_name, 'w', newline='') as csvfile:
        writer = csv.writer(csvfile)
        writer.writerows(darkness_list_background_subtracted)

    #export for making graphs
    my_model.dens_list = darkness_list_background_subtracted
    my_model.save()

    graph_types = ast.literal_eval(my_model.graph_types)
    strip_request = ast.literal_eval(my_model.strip_request)

    test_strips = stripper2(
        my_model.img_test_crop,graph_types,rows,cols,
        strip_request
        )
    mock_strips = stripper2(
        my_model.img_con_crop,graph_types,rows,cols,strip_request
        )

    my_dir = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
    my_dir = os.path.join(my_dir,my_model.user_id) 

    #save strips to folder for pptx
    for i in range(len(test_strips)):
        my_filename = 'strip_' + str(i)+'.png'
        test_strips[i].save(os.path.join(my_dir,my_filename)) 

    for i in range(len(mock_strips)):
        my_filename = 'mock_strip_' + str(i)+'.png'
        mock_strips[i].save(os.path.join(my_dir,my_filename)) 


    if my_string[last_slash_index + 1:] !='EMPTY.jpg':
        vis_strips = stripper2(
            my_model.img_vis_crop,graph_types,rows,cols,
            strip_request
            )
        my_dir = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
        my_dir = os.path.join(my_dir,my_model.user_id) 
        #save strips to folder for pptx
        for i in range(len(vis_strips)):
            my_filename = 'vis_strip_' + str(i)+'.png'
            vis_strips[i].save(os.path.join(my_dir,my_filename)) 


    #quick patch for to accomodate disc_y_lim
    graph_types_dummy = graph_types
    if 'd' in graph_types:
        graph_types_dummy = []
        for i in range(len(strip_request)):
            graph_types_dummy.append('h')

    #does vis exist for pptx strip adding
    if my_string[last_slash_index + 1:] !='EMPTY.jpg':  
        vis_bool = True
    else:
        vis_bool = False

    #make slides
    make_pptx(
        test_strips,graph_types_dummy,strip_request,my_dir,my_model.user_id,vis_bool
    )


    ### GRAPHS ###
    dens_list = my_model.dens_list
    #convert dens_list which is in array format to flat list of intensities to match strip_request indices
    dens_list_flat = []
    for col in range(cols):
        for row in range(rows):
            dens_list_flat.append(dens_list[row][col])

    #make graphs

    #disc array - calculate mx value from selected slice darknesses for y axis
    if 'd' in graph_types:
        #function to round max to nearest number of choice
        def myround20(x, base=20):
            return base * round((x+1)/base)

        my_flat_list = []
        for i in range(len(strip_request)):
            #convert strip request into slice indices
            new_strip = []
            new_strip.append(strip_request[i][0]-1)
            new_strip.append(strip_request[i][1])
            for i in dens_list_flat[new_strip[0]:new_strip[1]]:
                my_flat_list.append(i)
        disc_y_max = myround20(max_(my_flat_list))

        #return to v so graphs are made
        graph_types = []
        for i in range(len(strip_request)):
            graph_types.append('h')
    else:
        disc_y_max = False

    for i in range(len(strip_request)):
        #convert strip request into slice indices
        new_strip = []
        new_strip.append(strip_request[i][0]-1)
        new_strip.append(strip_request[i][1])

        my_heights = dens_list_flat[new_strip[0]:new_strip[1]]

        my_x = [j for j in range(len(my_heights))]
        my_filename = 'graph_'+str(i)
        if graph_types[i] == 'h':
            h_bar(my_x,my_heights,my_dir,my_filename,disc_y_max)
        if graph_types[i] == 'v':
            v_bar(my_x,my_heights,my_dir,my_filename,disc_y_max)


    #open ppt
    my_dir = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
    my_dir = os.path.join(my_dir,my_model.user_id) 
    my_filename = 'presentation_'+my_model.user_id+'.pptx'
    my_dir = os.path.join(my_dir,my_filename)

    #add graphs to ppt
    if my_dir == 'blank':
        prs = Presentation()
        prs.slide_width = Cm(33.9)
        prs.slide_height = Cm(19.1)
        for i in range(len(strip_request)):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
    else:
        prs = Presentation(my_dir)
    
    h_left = Cm(8.7)
    h_top = Cm(2.2)
    v_left = Cm(8.4)
    v_top = Cm(6.5)
    #width = Cm(0.7)
    
    slides_list = [slide.slide_id for slide in prs.slides]

    for i in range(len(strip_request)):
        slide = prs.slides.get(slides_list[i])

        #get path to graph
        my_dir = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
        my_dir = os.path.join(my_dir,my_model.user_id) 
        my_filename = "graph_"+str(i)+".png"
        my_dir = os.path.join(my_dir,my_filename)

        if graph_types[i] == "h":
            img = slide.shapes.add_picture(my_dir,h_left,h_top)
        if graph_types[i] == "v":
            img = slide.shapes.add_picture(my_dir,v_left,v_top)

    #save ppt
    my_dir = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
    my_dir = os.path.join(my_dir,my_model.user_id) 
    my_filename = 'presentation_'+my_model.user_id+'.pptx'
    my_dir = os.path.join(my_dir,my_filename)
    prs.save(my_dir)

    #compress files for output
    files_to_compress = []
    compression_filenames = []
    my_dir = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
    my_dir = os.path.join(my_dir,my_model.user_id) 
    for i in range(len(strip_request)):

        my_filename = 'graph_'+str(i)+'.png'
        files_to_compress.append(os.path.join(my_dir,my_filename))
        compression_filenames.append(my_filename)

        my_filename = 'strip_'+str(i)+'.png'
        files_to_compress.append(os.path.join(my_dir,my_filename))
        compression_filenames.append(my_filename)

        my_filename = 'mock_strip_'+str(i)+'.png'
        files_to_compress.append(os.path.join(my_dir,my_filename))
        compression_filenames.append(my_filename)

        my_string = my_model.img_vis.name
        last_slash_index = my_string.rindex('/')
        if my_string[last_slash_index + 1:] !='EMPTY.jpg':
            my_filename = 'vis_strip_'+str(i)+'.png'
            files_to_compress.append(os.path.join(my_dir,my_filename))
            compression_filenames.append(my_filename)

        files_to_compress.append(os.path.join(my_dir,'img_test.csv'))
        compression_filenames.append('img_test.csv')
        files_to_compress.append(os.path.join(my_dir,'img_con.csv'))
        compression_filenames.append('img_con.csv')
        files_to_compress.append(os.path.join(my_dir,'background_subtracted_data.csv'))
        compression_filenames.append('background_subtracted_data.csv')

        my_filename = 'presentation_'+my_model.user_id+'.pptx'
        files_to_compress.append(os.path.join(my_dir,my_filename))
        compression_filenames.append('DATA POWERPOINT.pptx')

    file_compress(files_to_compress,compression_filenames,os.path.join(my_dir,'data_peptide_array.zip'))


    return render(request, 'app/download.html', {})





def file_download(request):
    my_model = gridModel.objects.all().order_by('-id')[0]
    my_dir = os.path.join(MEDIA_ROOT,'sd') #!!!!!!!!!!!!!! Python anywhere problem
    my_dir = os.path.join(my_dir,my_model.user_id)
    file_path = os.path.join(my_dir,'data_peptide_array.zip')

    if os.path.exists(file_path): 

        with open(file_path, 'rb') as fh:
            response = HttpResponse(fh.read(), content_type="application/default")
            response['Content-Disposition'] = 'inline; filename=' + os.path.basename(file_path)
            return response

    raise Http404







def ss_form(request):
    return render(request, 'app/ss_form.html', {})

def ss_result(request):
    res=''
    seq = request.POST.get('enter-sequence').upper()
    m_notation_box = request.POST.get('m-notation')

    ala_box = request.POST.get('ala-scan')
    if ala_box == 'true':
        res += ala_scan(seq)
        res += '\n'
        if m_notation_box == 'true':
            res += '.space'
            res += '\n'
            res += '\n'

    n_trunc_box = request.POST.get('n-trunc')
    if n_trunc_box == 'true':
        res += n_trunc(seq)
        res += '\n'
        if m_notation_box == 'true':
            res += '.space'
            res += '\n'
            res += '\n'

    c_trunc_box = request.POST.get('c-trunc')
    if c_trunc_box == 'true':
        res += c_trunc(seq)
        res += '\n'
        if m_notation_box == 'true':
            res += '.space'
            res += '\n'
            res += '\n'

    n_c_trunc_box = request.POST.get('n-c-trunc')
    if n_c_trunc_box == 'true':
        res += n_c_trunc(seq)
        res += '\n'
        if m_notation_box == 'true':
            res += '.space'
            res += '\n'
            res += '\n'

    point_sub_box = request.POST.get('point-sub')
    if point_sub_box == 'true':
        res += point_sub(seq)
        res += '\n'
        if m_notation_box == 'true':
            res += '.space'
            res += '\n'
            res += '\n'
    return render(request, "app/ss_result.html", {"result":res})





def documentation(request):
    return render(request, 'app/documentation.html',{})

def tutorial(request):
    return render(request, 'app/tutorial.html',{})